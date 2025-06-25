import re
import PyPDF2
import io
import numpy as np
import easyocr
from PIL import Image, ImageDraw
from pdf2image import convert_from_path
import docx2txt
import PyPDF2
from pptx import Presentation

def convert_text_to_words(text):
    # Use regular expressions to split the text into words
    words = re.findall(r'\b\w+\b', text)
    return words

def detect_document_type(text):
    regulatory_words=['archive', 'connectivity', 'security', 'distribution', 'conviction', 'reasoning', 'damage', 'substance', 'cross-contamination', 'implementation', 'indemnity', 'development', 'comity', 'incentive', 'outbreak', 'admission', 'allegation', 'well-being', 'incident', 'file', 'investigation', 'vector', 'discretion', 'code', 'disbursement', 'usage', 'certification', 'authorization', 'presumption', 'withdrawal', 'protocol', 'collaboration', 'requirement', 'service', 'random', 'blocking', 'range', 'credentials', 'abandonment', 'high-risk', 'feedback', 'written', 'consultation', 'consent', 'correspondence', 'causality', 'commencement', 'substitution', 'assessment', 'declaration', 'conclusion', 'uncertainty', 'closure', 'prediction', 'permit', 'specification', 'decontamination', 'reception', 'offense', 'strategy', 'environmental', 'infraction', 'indication', 'condition', 'licensing', 'promotion', 'accuracy', 'autonomy', 'solution', 'validation', 'regulation', 'classification', 'ambiguity', 'auditability', 'consequence', 'case', 'communication', 'supply', 'structure', 'disclosure', 'discovery', 'capability', 'backup', 'containment', 'precision', 'consistency', 'manipulation', 'shield', 'blueprint', 'manual', 'exclusion', 'performance', 'threat', 'duty', 'fumigation', 'court', 'evaluation', 'concurrency', 'derivation', 'cross-reference', 'deviation', 'bond', 'protection', 'transparency', 'cleanliness', 'accusation', 'conduct', 'forecast', 'readiness', 'disruption', 'acquisition', 'abatement', 'guidelines', 'treatment', 'intervention', 'response', 'execution', 'concealment', 'exposure', 'discipline', 'liability', 'outage', 'delusion', 'enforcement', 'interpretation', 'filling', 'support', 'reinforcement', 'denial', 'ignition', 'conformity', 'key', 'pain', 'action', 'organization', 'posture', 'quarantine', 'strengthening', 'deduction', 'fire', 'determination', 'breach', 'amendment', 'dedication', 'impression', 'import', 'trigger', 'contravention', 'data', 'apparatus', 'event', 'capacity', 'injection', 'biohazard', 'privacy', 'result', 'campaign', 'adoption', 'collation', 'emergency', 'accident', 'limitation', 'selection', 'workforce', 'error', 'deterrence', 'clauses', 'sensitivity', 'initiative', 'collection', 'fines', 'hygiene', 'biosecurity', 'alert', 'clarity', 'issuance', 'actuation', 'confidentiality', 'finalization', 'procurement', 'litigation', 'safety', 'criticality', 'contract', 'commitment', 'binding', 'log', 'dismissal', 'information', 'recovery', 'inference', 'customization', 'release', 'skill', 'benchmark', 'countermeasure', 'first aid', 'expertise', 'qualification', 'announcement', 'openness', 'method', 'provision', 'production', 'heat', 'integration', 'segregation', 'delivery', 'management', 'tolerance', 'procedure', 'crime', 'avoidance', 'product', 'control', 'capture', 'demarcation', 'retention', 'articulation', 'circulation', 'buffer', 'awareness', 'trust', 'governance', 'conversion', 'prohibition', 'infection', 'verification', 'flow', 'quantity', 'vibration', 'consumer', 'analysis', 'biometric', 'cybersecurity', 'insurance', 'software', 'access', 'collateral', 'reimbursement', 'independence', 'false', 'temperature', 'accreditation', 'objectivity', 'variance', 'fraud', 'violation', 'confession', 'restraint', 'confirmation', 'arbitration', 'conformance', 'claim', 'fitness', 'maintenance', 'warning', 'discrimination', 'sterilization', 'cancellation', 'water', 'capabilities', 'mapping', 'conservation', 'integrity', 'dependency', 'packaging', 'hazard', 'workplace', 'surveillance', 'appraisal', 'review', 'follow-up', 'recognition', 'currency', 'pharmaceutical', 'justification', 'penalty', 'removal', 'efficacy', 'instructions', 'laboratory', 'characterization', 'redress', 'audit', 'exemption', 'documentation', 'generalization', 'lockout', 'notification', 'omission', 'infringement', 'unauthorized', 'continuity', 'personnel', 'planning', 'stability', 'connection', 'technical', 'ban', 'facility', 'consideration', 'inspection', 'calibration', 'discontinuity', 'configuration', 'sanction', 'clarification', 'collective', 'resilience', 'diversion', 'preparation', 'cybernetics', 'training', 'signature', 'illustration', 'correlation', 'eligibility', 'simulation', 'source', 'storage', 'directive', 'restriction', 'quality', 'coordination', 'sanitization', 'examination', 'prevention', 'defense', 'assignment', 'risk', 'infrastructure', 'dispute', 'subcontractor', 'report', 'telecommunication', 'deadline', 'adherence', 'weakness', 'scheme', 'diligence', 'propagation', 'recall', 'vehicle', 'system', 'dissolution', 'target', 'standard', 'panel', 'compliance', 'remediation', 'benefit', 'monitoring', 'allocation', 'update', 'radiation', 'submission', 'waste', 'ethics', 'medication', 'convenience', 'computation', 'challenge', 'congruence', 'defect', 'disposal', 'formality', 'change', 'confusion', 'constancy', 'establishment', 'pesticide', 'coherence', 'web', 'corrective', 'remittance', 'default', 'resolution', 'inclusion', 'marking', 'feasibility', 'improper', 'suitability', 'assertion', 'deployment', 'suspension', 'frequency', 'fault', 'consensus', 'designation', 'pollution', 'self-inspection', 'document', 'help', 'identification', 'data integrity', 'advisory', 'failure', 'transport', 'time', 'diagnosis', 'oversee', 'consolidation', 'infiltration', 'material', 'reliability', 'escalation', 'discontinuation', 'policy', 'alteration', 'affiliation', 'disposition', 'inquiry', 'attestation', 'optimization', 'interference', 'crisis', 'coverage', 'correction', 'observation', 'extension', 'definition', 'succession', 'paramedic', 'intensity', 'severity', 'measurement', 'public', 'mitigation', 'gap', 'sustainability', 'anticipation', 'dictate', 'authentication', 'testing', 'labeling', 'assurance', 'screening', 'operation', 'adverse', 'preemption', 'grievance', 'nullification', 'accountability', 'approval', 'concession', 'settlement', 'corroboration', 'demand', 'behavior', 'grant', 'precaution', 'detection', 'cooperation', 'categorization', 'portable']
    financial_words =['call', 'portfolio diversification', 'security', 'ingredient', 'business', 'active', 'market', 'spread', 'advance', 'depreciation', 'payment', 'federal', 'indemnity', 'development', 'global', 'student loan', 'return', 'trademark', 'dow', 'will', 'decision', 'preclinical', '401(k)', 'rule of 72', 'growth stock', 'revenue stream', 'pharmacovigilance', 'supplier', 'investigation', 'dosage form', 'inflation', 'capital gain', 'net income', 'commodity', 'disbursement', 'innovation', 'debt', 'withdrawal', 'collaboration', 'cash flow', 'yield curve', 'basis', 'outcomes', 'over-the-counter (otc)', 'cro (contract research organization)', 'divestiture', 'preapproval', 'insurance agent', 'zero-coupon bond', 'delisting', 'convertible', 'shareholder', 'rent', 'restructuring', 'post-marketing', 'divestment', 'sales', 'intellectual property', 'payroll clerk', 'strategy', 'municipal bond', 'term life insurance', 'subsidy', 'marketing', 'blue chip', 'licensing', 'off-label', 'dossier', 'foreclosure', 'roth', 'regulation', 'credit score', 'value', 'orphan', 'target market', 'direct deposit', 'supply', 'demerger', 'subsidiary', 'microfinance', 'appreciation', 'nasdaq', 'actuary', 'disclosure', 'firm', 'drug', 'automated teller machine (atm)', 'nyse (new york stock exchange)', 'purchase', 'certificate', 'tax exemption', 'expiration', 'consultant', 'discounting', 'customer', 'niche', 'money market account', 'advisor', 'capital', 'tax deduction', 'underwater mortgage', 'sec (securities and exchange commission)', 'bond', 'acquire', 'average', 'operating', 'dollar', 'bull market', 'forecast', 'royalty', 'mortgage', 'acquisition', 'invest', 'payroll tax', 'liability', 'franchise', 'net', 'rebranding', 'earn', 'support', 'action', 'secured loan', 'growth', 'collateralized debt obligation (cdo)', 'universal life insurance', 'trader', 'import', 'payer', 'bank', 'data', 'liquidity', 'risk manager', 'hospital', 'sole proprietorship', 'event', 'oligopoly', 'job', 'financial literacy', 'economy', 'clearing', 'decentralization', 'workforce', 'fee', 'penny stock', 'quality control', 'voluntary bankruptcy', 'wholesale', 'charge', 'loan officer', 'regulatory', 'biosecurity', 'arbitrage', 'fda approval', 'export', 'confidentiality', 'procurement', 'mutual', 'litigation', 'exchange', 'bondholder', 'retirement account', 'contract', 'commitment', 'medical', 'technology', 'costing', 'globalization', 'customization', 'property', 'broker', 'debenture', 'benchmark', 'discount', 'gain', 'line of credit', 'therapeutic', 'overdraft', 'deposit', 'mortgage rate', 'corporate', 'alliance', 'recession', 'unemployment', 'tax', 'annuity', 'share', 'hedge', 'investment portfolio', 'price', 'placebo', 'management', 'expense', 'equity', 'annual', 'certificate of deposit (cd)', 'commercial', 'short selling', 'credit card', 'product', 'control', 'dilution', 'commodities market', 'asset', 'compound', 'insurance policy', 'venture', 'trust', 'governance', 'counterfeit', 'public offering', 'diversification', 'buy', 'rate', 'merger', 'analysis', 'preferred stock', 'private equity', 'leverage ratio', 'agency', 'pharmacy', 'insurance', 'access', 'liquid asset', 'collateral', 'patent', 'bioequivalence', 'lease', 'estate', 'fraud', 'order', 'underwrite', 'p/e (price-to-earnings) ratio', 'lender', 'claim', 'balance sheet', 'gdp (gross domestic product)', 'simple interest', 'financial advisor', 'revenue', 'trade', 'treasury bill', 'treasury', 'biosafety', 'maturity', 'savings account', 'initial deposit', 'borrow', 'cost', 'entrepreneur', 'stockbroker', 'packaging', 'authority', 'agreement', 'savings', 'stock', 'accountant', 'index', 'accrual', 'option', 'futures', 'economist', 'chartered financial analyst (cfa)', 'cashier', 'currency', 'adjusted', 'roth ira', 'principal', 'institutional', 'annual report', 'online banking', 'efficacy', 'payable', 'laboratory', 'audit', 'close', 'dividend', 'inventory', 'merchant', 'dividend yield', 'tax return', 'cash', 'minimum payment', 'apr (annual percentage rate)', 'foreign', 'research', 'refinance', 'auditor', 'supply chain', 'facility', 'finance', 'bankruptcy', 'joint account', 'volatility', 'transaction', 'compound interest', 'underwriter', 'alternative investment', 'financial planner', 'deductible', 'biosimilar', 'amortization schedule', 'whole life insurance', 'expenditure', 'term', 'clinical', 'loss', 'quality', 'r&d (research and development)', 'fixed', 'penny', 'ipo (initial public offering)', 'fixed income', 'pricing strategy', 'dow jones industrial average (djia)', 'bull', 'yield', 'spac (special purpose acquisition company)', 'biologic', 'withholding tax', 'default risk', 'budget', 'brand', 'derivative', 'risk', 'prime rate', 'exclusivity', 'write-off', 'spot market', 'deficit', 'employee benefits', 'initial', 'recall', 'account', 'leverage', 'sales revenue', 'quote', 'commercial bank', 'standard', 'real', 'compliance', 'benefit', 'monitoring', 'allocation', 'time value of money', 'gross income', 'financial', 'profit', 'overdraft protection', 'price elasticity', 'medicare', 'account balance', 'banker', 'swap', 'prescription', 'analyst', 'emergency fund', 'union', 'partnership', 'corporate bond', 'investment', 'balance', 'default', 'buy-and-hold', 'federal reserve system (the fed)', 'junk bond', 'cash value', 'pipeline', 'day trading', 'fda (food and drug administration)', 'vendor', 'sale', 'counterfeiting', 'loan', 'economic', 'fund', 'government', 'ytd (year-to-date)', 'bear market', 'real estate investment trust (reit)', 'capitalization', 'company', 'amortization', 'pension', 'installment loan', 'free cash flow', 'credit', 'short', 'retained earnings', 'futures contract', 'over-the-counter (otc)', 'dosage', 'consolidation', 'price fixing', 'earnings', 'efficiency', 'investor', 'portfolio', 'segmentation', 'policy', 'mortgage broker', 'budget deficit', 'interest rate', 'bookkeeper', 'merchant account', 'limit', 'variable interest rate', 'wage garnishment', 'exchange rate', 'offer', 'crisis', 'appraiser', 'fund manager', 'home equity', 'securities', 'interest', 'technology transfer', 'employee', 'option contract', 'call option', 'public', 'market capitalization', 'withhold', 'adjustable-rate mortgage', 'wealth', 'in-licensing', 'medicaid', 'stock exchange', 'formulation', 'cryptocurrency', 'usury', 'generic', 'refinancing', 'inflation rate', 'price skimming', 'income']
    medical_words=['medical','diagnosed','treatment','health','viral', 'nephritis', 'birth', 'colitis', 'discomfort', 'aortic', 'acid', 'erectile', 'integumentary', 'basal', 'mastectomy', 'fistula', 'hyperlipidemia', 'echocardiography', 'detoxification', 'astigmatism', 'splenomegaly', 'morbidity', 'addiction', 'abnormal', 'thyroid', 'arthritis', 'rheumatoid', 'fluorescence', 'dystrophy', 'metabolism', 'sinusitis', 'eczema', 'hypertension', 'cell', 'genital', 'obesity', 'hyperactivity', 'embolism', 'dental', 'monitor', 'lymphoma', 'sore', 'mastitis', 'heart', 'cornea', 'arterial', 'electrocardiogram', 'localized', 'analgesic', 'insomnia', 'carcinoma', 'vaccine', 'spasm', 'electroencephalogram', 'abdominal', 'vision', 'ear', 'incontinence', 'lithium', 'glucose', 'muscle', 'medullary', 'angioplasty', 'fibrosis', 'anterior', 'gastrointestinal', 'infiltrate', 'inflammatory', 'varicose', 'cognition', 'trauma', 'esophagitis', 'gallstone', 'antidepressant', 'drug', 'acupuncture', 'inhibitor', 'denture', 'neutropenia', 'expiration', 'hepatic', 'convalescence', 'allergy', 'molar', 'interstitial', 'gastroenterology', 'alternative', 'thrombosis', 'nutrition', 'brain', 'mitral', 'psychiatry', 'ascites', 'aorta', 'fasting', 'deviation', 'electrocardiography', 'cough', 'intravenous', 'conjunctiva', 'inhalation', 'yeast', 'miosis', 'exercise', 'endocrine', 'catheter', 'headache', 'infectious', 'goiter', 'lipid', 'hypocalcemia', 'osteoarthritis', 'diverticulitis', 'lung', 'diverticulosis', 'intervention', 'mucus', 'antibiotic', 'urinary', 'metabolic', 'hypoglycemia', 'dermatology', 'hallucination', 'pain', 'enema', 'mri', 'lesion', 'cerebellum', 'polyp', 'growth', 'neurology', 'radiology', 'esophageal', 'congenital', 'hospital', 'cartilage', 'endoscope', 'metabolite', 'injection', "alzheimer's", 'lymphatic', 'disorder', 'pediatrics', 'hernia', 'essential', 'emergency', 'palliative', 'cirrhosis', 'gastroenteritis', 'memory', 'epidural', 'anticoagulant', 'candidiasis', 'avascular', 'speech', 'tuberculosis', 'melanoma', 'anesthetic', 'immunodeficiency', 'lumen', 'psoriasis', 'contracture', 'sleep', 'cellulitis', 'whiplash', 'sexually', 'glomerulus', 'cystic', 'intubation', 'constipation', 'enzyme', 'hematology', 'aneurysm', 'cleft', 'copd', 'aggression', 'atherosclerosis', 'dissection', 'tachycardia', 'brachytherapy', 'kidney', 'granuloma', 'bloodstream', 'glaucoma', 'hip', 'skeletal', 'hemorrhoids', 'guillain-barre', 'gangrene', 'idiopathic', 'acute', 'hyperventilation', 'cardiac', 'keratitis', 'dose', 'epithelium', 'wellness', 'coughing', 'gallbladder', 'respiratory', 'hormonal', 'infection', 'immunization', 'tumor', "parkinson's", 'cataract', 'erythrocyte', 'depression', 'hypotension', 'intramuscular', 'gastritis', 'impairment', 'bariatric', 'antenatal', 'menopause', 'epilepsy', 'blood', 'diffusion', 'cochlear', 'vein', 'osteoporosis', 'anxiety', 'bradycardia', 'depressant', 'hepatitis', 'mammography', 'rehabilitation', 'vasectomy', 'cranial', 'dementia', 'pelvic', 'endoscopy', 'hyperplasia', 'cardiovascular', 'bone', 'tendinitis', 'heartburn', 'surgery', 'clotting', 'limb', 'fracture', 'conjunctivitis', 'microbiology', 'emphysema', 'exfoliation', 'colonoscopy', 'hematemesis', 'carpal', 'excision', 'infertility', 'migraine', 'immune', 'hypervolemia', 'halitosis', 'dysphagia', 'measles', 'comorbidity', 'bronchitis', 'intensive', 'cervical', 'breathing', 'bursitis', 'inpatient', 'hormone', 'bladder', 'ligament', 'luminal', 'calcium', 'efficacy', 'seizure', 'laboratory', 'sprain', 'convulsion', 'cognitive', 'coronary', 'dilatation', 'pregnancy', 'corticosteroid', 'autoimmune', 'folate', 'diarrhea', 'etiology', 'meatus', 'cerebrospinal', 'optometry', 'digestion', 'mineral', 'geriatric', 'hemiparesis', 'infarction', 'ileus', 'alcoholism', 'leukemia', 'fever', 'earache', 'symptom', 'ectopic', 'ambulatory', 'laryngitis', 'mumps', 'inflammation', 'aspiration', 'clinical', 'outpatient', 'reflux', 'laxative', 'cardiology', 'hereditary', 'influenza', 'mania', 'circulatory', 'impaction', 'examination', 'angina', 'cystectomy', 'biopsy', 'hospice', 'coagulation', 'cancer', 'overdose', "behcet's", 'ulcer', 'stomach', 'weight', 'analgesia', 'wound', 'stroke', 'pathology', 'transplant', 'x-ray', 'gout', 'colorectal', 'anesthesia', 'malnutrition', 'induration', 'menstruation', 'diagnostic', 'diabetes', 'methadone', 'cannula', 'dehydration', 'pulmonary', 'spinal', 'digestive', 'nausea', 'wart', 'ophthalmology', 'crepitus', 'herniation', 'intraventricular', 'dermatitis', 'immunoglobulin', 'adenocarcinoma', 'flatulence', 'excretion', 'medication', 'lumpectomy', 'edema', 'clavicle', 'ataxia', 'cholecystectomy', 'syndrome', 'aphasia', 'rashes', 'distension', 'dentistry', 'liver', 'adrenal', 'fatigue', 'syphilis', 'malignant', 'dizziness', 'anorexia', 'anaphylaxis', 'leukocyte', 'anemia', 'agitation', 'bruise', 'appendectomy', 'cholesterol', 'discharge', 'hematoma', 'atrophy', 'bunion', 'electrolyte', 'median', 'affective', 'hypothermia', 'skin', 'tonsillitis', 'chemotherapy', 'multiple', 'pneumonia', 'disability', 'bile', 'antiseptic', 'arrhythmia', 'mammogram', 'noninvasive', 'urology', 'genetics', 'carotid', 'sepsis', 'epistaxis', 'electroconvulsive', 'gastric', 'malignancy', 'bioavailability', 'lymphocyte', 'ganglion', 'vomiting', 'antipsychotic', 'laryngeal', 'jaundice', 'diagnosis', 'dosage', 'feces', 'asthma', 'endocrinology', "barrett's", 'ekg', 'malaria', 'pancreatic', 'immunity', 'miscarriage', 'extracorporeal', 'alertness', 'impotence', 'coma', 'ischemia', 'auditory', 'body', 'hemorrhage', 'infusion', 'hypothyroidism', 'renal', 'fibromyalgia', 'podiatry', 'eye', 'axon', 'endoscopic', 'abdomen', 'abscess', 'cold', 'diuretic', 'swelling', 'bowel', 'microscope', 'iatrogenic', 'diet', 'prostate', 'celiac', 'carbohydrate', 'dietetics', 'senility', 'chronic', 'gynecology', 'joint', 'benign', 'autism', 'balloon', 'lactation', 'cerebral', 'meningitis']
    legal_words=['hearsay', 'extortion', 'vested', 'infringe', 'refinance', 'imprisonment', 'judge', 'contempt', 'forfeiture', 'acquittal', 'legitimate', 'retire', 'mandate', 'lienholder', 'notice', 'bankrupt', 'domicile', 'rescission', 'expunge', 'defendant', 'ordinance', 'prohibit', 'enforceable', 'prosecution', 'malfeasance', 'compliance', 'bond', 'counsel', 'predecessor', 'evidence', 'forcible', 'incarcerate', 'affidavit', 'decree', 'misrepresent', 'subrogate', 'noncompliance', 'offensive', 'infirmity', 'cite', 'contingency', 'enjoin', 'lawsuit', 'abatement', 'coercion', 'nonwaivable', 'surety', 'bequest', 'merger', 'sequestration', 'treason', 'leasehold', 'grievance', 'guarantee', 'capital', 'amicus', 'fiduciary', 'caveat', 'abeyance', 'wrongful', 'annulment', 'constitutional', 'rebut', 'advocate', 'pardon', 'return', 'reversal', 'intruder', 'moot', 'misdemeanor', 'affiant', 'easement', 'transient', 'nominee', 'testimony', 'inaction', 'vouch', 'appealable', 'anticipatory', 'secure', 'brief', 'certiorari', 'omission', 'entitle', 'relinquishment', 'mortgage', 'unlawful', 'demurrer', 'constitutionality', 'judicial', 'mitigation', 'credence', 'exemption', 'construe', 'irrevocable', 'ex parte', 'litigant', 'judgment', 'inhibit', 'subpoena', 'ex post facto', 'intestate', 'swindle', 'conveyance', 'witness', 'restitution', 'inception', 'protect', 'threat', 'remunerate', 'inhabitant', 'interrogate', 'negotiable', 'sentencing', 'defer', 'usury', 'eviction', 'prima', 'appellee', 'contract', 'docket', 'defamation', 'deed', 'chattel', 'conceal', 'intrusion', 'transitory', 'impeachment', 'nonjoinder', 'answer', 'portable', 'quantum', 'testimonial', 'restrict', 'exempt', 'notify', 'execution', 'detriment', 'reasonable', 'referee', 'affirmation', 'possession', 'larceny', 'hereof', 'interim', 'warrant', 'pursuant', 'party', 'illegality', 'promote', 'prejudice', 'nonresident', 'section', 'insolvent', 'substitution', 'creditor', 'lease', 'arraignment', 'pretrial', 'moral', 'injury', 'confession', 'conviction', 'mandator', 'remember', 'restoration', 'hereinafter', 'defend', 'proceedings', 'process', 'citation', 'continuance', 'de jure', 'intervene', 'solemn', 'successor', 'undertake', 'bill', 'corporate', 'compromise', 'parol', 'pleading', 'lessor', 'judiciary', 'subordinate', 'interpose', 'curtilage', 'caution', 'incarceration', 'bail', 'nonfeasance', 'mitigating', 'damages', 'summons', 'tribunal', 'law', 'lenient', 'jurisprudence', 'execute', 'remand', 'lis', 'retroactive', 'transgression', 'counterclaim', 'victimize', 'proxy', 'appointment', 'vicarious', 'attorney', 'vexatious', 'surrender', 'indorse', 'derogate', 'taxpayer', 'solicit', 'contractor', 'ostensible', 'legacy', 'condemn', 'mediation', 'leniency', 'penalty', 'preclosure', 'supervision', 'unjust', 'sovereign', 'unenforceable', 'fugitive', 'appraise', 'preemptive', 'misdeed', 'termination', 'devisee', 'petition', 'manslaughter', 'speculation', 'legislation', 'indemnity', 'notwithstanding', 'attest', 'appraisal', 'franchise', 'perjury', 'incorporeal', 'facilitate', 'appeal', 'reserve', 'precedent', 'resolve', 'expert', 'impute', 'rights', 'repudiate', 'proffer', 'usurpation', 'suspicion', 'jeopardize', 'intention', 'interrogatories', 'counterfeit', 'impeach', 'premature', 'reimburse', 'forbid', 'lapse', 'appellate', 'representative', 'immaterial', 'intimidate', 'offense', 'tortfeasor', 'punitive', 'obligation', 'clause', 'testify', 'adjournment', 'indictment', 'entitlement', 'criminal', 'equitable', 'rebuttable', 'liable', 'jurisdiction', 'ingress', 'transgress', 'negotiate', 'renounce', 'legal', 'violation', 'writ', 'garnishment', 'homicide', 'plea', 'rule', 'order', 'adversary', 'ordinarily', 'statement', 'tenancy', 'bench', 'perpetrator', 'emolument', 'injunction', 'revocation', 'devise', 'mitigate', 'prospective', 'alibi', 'indemnify', 'disclaimer', 'amendment', 'attorney-in-fact', 'exculpate', 'negotiator', 'reprieve', 'undisclosed', 'deficiency', 'sue', 'proceeding', 'represent', 'solicitor', 'exonerate', 'delegable', 'ejectment', 'decision', 'quorum', 'dissent', 'cross-examination', 'stipulate', 'consent', 'waiver', 'bequeath', 'obligate', 'lawful', 'fraud', 'arbitrage', 'vandalism', 'privilege', 'dictum', 'mense', 'dissolution', 'retrial', 'rescind', 'unreasonable', 'prison', 'plead', 'quash', 'parliament', 'mandamus', 'eminent', 'estoppel', 'statute', 'estray', 'immunity', 'collateral', 'substantiate', 'hostile', 'lessee', 'publish', 'portability', 'cajole', 'claimant', 'innocent', 'court', 'investigate', 'commission', 'refugee', 'paralegal', 'forgery', 'mediator', 'disbarment', 'complaint', 'slander', 'reformation', 'reservation', 'right', 'de facto', 'executor', 'conflict', 'sponsor', 'corroborate', 'disclose', 'bestow', 'disorder', 'restrain', 'premises', 'lord', 'exclusionary', 'interlocutory', 'peremptory', 'summon', 'plaintiff', 'inadmissible', 'statutory', 'treasurer', 'victim', 'violence', 'chambers', 'violate', 'abrogate', 'bailment', 'lien', 'forfeit', 'exculpatory', 'dismissal', 'certification', 'conform', 'instigate', 'barrister', 'lawyer', 'transit', 'marital', 'donor', 'warranty', 'sentence', 'verdict', 'accused', 'supplementary', 'settlement', 'convict', 'justice', 'issue', 'stake', 'arrest', 'nonperformance', 'embezzle', 'confer', 'innocuous', 'appellant', 'demise', 'extradition', 'sequester', 'preponderance', 'rational', 'presumptuous', 'debtor', 'garnish', 'situation', 'premise', 'profession', 'dismiss', 'forensic', 'enforceability', 'prosecutor', 'equity', 'invent', 'direct', 'hinder', 'power', 'guilty', 'patent', 'recourse', 'breach', 'conspire', 'monetary', 'inducement', 'preservation', 'proclamation', 'voluntary', 'incorporate', 'consortium', 'intervenor', 'conclusive', 'title', 'negotiation', 'testament', 'alienation', 'circumstantial', 'rebuttal', 'oath', 'deprecate', 'foreclosure', 'withhold', 'arbitration', 'prudent', 'waive', 'proof', 'seizure', 'deliberation', 'adjudication', 'consult', 'trial', 'prosecutable', 'precaution', 'presumption', 'induce', 'relinquish', 'infringement', 'strike', 'exhibit', 'prosecute', 'onus', 'denomination', 'assumpsit', 'misfeasance', 'supervise', 'nullification']

    
    words = convert_text_to_words(text)
    words = [word.lower() for word in words]

    word_count = len(words)

    legal_count = sum(word.lower() in legal_words for word in words)
    financial_count = sum(word.lower() in financial_words for word in words)
    medical_count = sum(word.lower() in medical_words for word in words)
    rqa_count = sum(word.lower() in regulatory_words for word in words)

    legal_percentage = (legal_count / word_count) * 100 if word_count != 0 else 0
    financial_percentage = (financial_count / word_count) * 100 if word_count != 0 else 0
    medical_percentage = (medical_count / word_count) * 100 if word_count != 0 else 0
    rqa_percentage = (rqa_count / word_count) * 100 if word_count != 0 else 0



    if word_count == 0:
        return 'Unknown Document Type'
    
    max_percent = max(legal_percentage, financial_percentage, medical_percentage,rqa_percentage)

    # threshold of 4% for words detected in each class
    
    if max_percent == legal_percentage and legal_percentage>=4:
        legal_words_detected = set([word for word in words if word.lower() in legal_words])
        print("Legal Words Detected: ", legal_words_detected)
        print("Probability of Legal Document: {:.2f}%".format(legal_percentage))
        return 'Legal Document'
    elif max_percent == financial_percentage and financial_percentage>=4:
        financial_words_detected = set([word for word in words if word.lower() in financial_words])
        print("Financial Words Detected: ", financial_words_detected)
        print("Probability of Financial Document: {:.2f}%".format(financial_percentage))    
        return 'Financial Document'
    elif max_percent == medical_percentage and medical_percentage>=4:
        medical_words_detected = set([word for word in words if word.lower() in medical_words])
        print("Medical Words Detected: ", medical_words_detected)
        print("Probability of Medical Document: {:.2f}%".format(medical_percentage))
        return 'Medical Document'
    elif max_percent == rqa_percentage and rqa_percentage>=4:
        rqa_words_detected = set([word for word in words if word.lower() in regulatory_words])
        print("Regulatory/Quality/Safety Words Detected: ", rqa_words_detected)
        print("Probability of Regulatory/Quality/Safety Document: {:.2f}%".format(rqa_percentage))
        return 'Regulatory/Quality/Safety Document'
    else:
        return 'Unknown Document Type'
    


# def extract_words_from_pdf(file):
#     words = []
    
#     with open(file_path, 'rb') as file:
#         reader = PyPDF2.PdfReader(file)
        
#         for page in reader.pages:
#             text = page.extract_text()
#             words.extend(text.split())
    
#     return words

#     words = []
#     with io.BytesIO(file.read()) as file_buffer:
#         reader = PyPDF2.PdfReader(file_buffer)

#         for page in reader.pages:
#             text = page.extract_text()
#             words.extend(text.split())

#     return words


def extract_words_from_document(file):
    words = []

    file_extension = file.name.split('.')[-1].lower()

    if file_extension == 'pdf':
        words.extend(extract_words_from_pdf(file))
    elif file_extension == 'docx':
        words.extend(extract_words_from_docx(file))
    elif file_extension == 'pptx':
        words.extend(extract_words_from_pptx(file))
    else:
        image_data = file.read()
        words.extend(extract_words_from_image(image_data))

    return words

def extract_words_from_pdf(file):
    words = []
    print(words)

    with io.BytesIO(file.read()) as file_buffer:
        reader = PyPDF2.PdfReader(file_buffer)

        for page in reader.pages:
            text = page.extract_text()
            words.extend(text.split())

    return words

def extract_words_from_docx(file):
    text = docx2txt.process(file)
    words = text.split()

    return words

def extract_words_from_pptx(file):
    words = []

    presentation = Presentation(file)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                words.extend(shape.text.split())

    return words

def extract_words_from_image(image_data):
    if hasattr(image_data, 'read'):  # Check if the input is a file object
        image_data = image_data.read()

    reader = easyocr.Reader(['en'])
    result = reader.readtext(image_data)

    words = [item[1] for item in result]

    return words